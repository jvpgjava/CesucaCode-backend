from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

SUPPORTED_EXTENSIONS = {"pdf", "docx", "pptx", "txt"}


class UnsupportedFileTypeError(Exception):
    pass


def get_extension(filename: str) -> str:
    if "." not in filename:
        return ""
    return filename.rsplit(".", 1)[-1].lower()


def extract_text(file, filename: str) -> str:
    ext = get_extension(filename)
    if ext == "pdf":
        return _extract_pdf(file)
    if ext == "docx":
        return _extract_docx(file)
    if ext == "pptx":
        return _extract_pptx(file)
    if ext == "txt":
        return _extract_txt(file)
    raise UnsupportedFileTypeError(f"Formato .{ext} não suportado.")


def _extract_pdf(file) -> str:
    reader = PdfReader(file)
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(p for p in pages if p.strip())


def _extract_docx(file) -> str:
    doc = DocxDocument(file)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_pptx(file) -> str:
    presentation = Presentation(file)
    slides = []
    for slide in presentation.slides:
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if texts:
            slides.append("\n".join(texts))
    return "\n\n".join(slides)


def _extract_txt(file) -> str:
    return file.read().decode("utf-8", errors="ignore")
